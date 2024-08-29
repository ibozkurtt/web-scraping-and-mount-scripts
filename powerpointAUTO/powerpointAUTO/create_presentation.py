import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

def resize_image(input_img_path, output_img_path, max_width, max_height):
    original_img = Image.open(input_img_path)
    
    # Orjinal boyutları al
    original_width, original_height = original_img.size

    # Hedef boyutları belirle
    new_width = min(original_width, max_width)
    new_height = min(original_height, max_height)

    # Yeni boyutlarına göre orantılı bir şekilde küçült
    resized_img = original_img.resize((new_width, new_height), Image.Resampling.LANCZOS)

    # Yeni boyutlara sahip resmi kaydet
    resized_img.save(output_img_path, format="JPEG", quality=95)

def create_presentation(input_images_folder, output_images_folder, save_path):
    prs = Presentation()

    input_images = [file for file in os.listdir(input_images_folder) if file.endswith(('jpg', 'jpeg', 'png'))]

    for input_image_name in input_images:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = f"Input: {input_image_name}"

        input_img_path = os.path.join(input_images_folder, input_image_name)
        pic = slide.shapes.add_picture(input_img_path, Inches(0), Inches(0), width=Inches(1.8), height=Inches(2.5))

        output_folder_name = f"Output_{input_image_name}"
        output_folder_path = os.path.join(output_images_folder, output_folder_name)

        if os.path.exists(output_folder_path):
            output_images = [output_image for output_image in os.listdir(output_folder_path) if not output_image.startswith('.DS_Store') and output_image.lower().endswith(('jpg', 'jpeg', 'png'))]

            left_output = Inches(0)
            top_output = Inches(2.5)  # Üst kısmın başlangıcı

            for i, output_image_name in enumerate(output_images):
                resized_output_img_path = os.path.join(output_folder_path, output_image_name)
                resize_image(os.path.join(output_folder_path, output_image_name), resized_output_img_path, max_width=Inches(5.08), max_height=Inches(5.715))  # Çıkış resimlerini boyutlandır
                pic_output = slide.shapes.add_picture(resized_output_img_path, left_output, top_output, width=Inches(2), height=Inches(2.5))
                
                # Her beş resimden sonra bir sonraki satıra geç
                if (i + 1) % 5 == 0:
                    left_output = Inches(0)
                    top_output += Inches(2.5)  # 1.143 inç = 2.9 cm (1 inç = 2.54 cm)
                else:
                    left_output += Inches(2)
        else:
            print(f"Warning: {output_folder_path} not found. Skipping.")

    prs.save(save_path)

if __name__ == "__main__":
    input_folder_path = "/Users/ibrahimbozkurt/Desktop/powerpointAUTO/choosen_Images"
    output_folder_path = "/Users/ibrahimbozkurt/Desktop/powerpointAUTO/Results"
    save_path = "/Users/ibrahimbozkurt/Desktop/powerpointAUTO/output_presentation.pptx"
    create_presentation(input_folder_path, output_folder_path, save_path)
